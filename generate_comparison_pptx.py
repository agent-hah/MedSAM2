"""
Generate a PowerPoint presentation comparing model predictions side-by-side.

For each of 3 comparison pairs (FSL, SSL, WSL), shows:
  Input | Model A Prediction | Model B Prediction
across 4 selected cases (2 rat, 2 human).

Presentation structure:
  Title → FSL divider → Rat sub-header → 2 rat slides → Human sub-header → 2 human slides
       → SSL divider → ...
       → WSL divider → ...

Usage:
    conda run -n medsam2 python generate_comparison_pptx.py
"""

import os
from pathlib import Path

import numpy as np
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------
RESULTS_DIR = Path(__file__).parent / "results"
OUTPUT_FILE = RESULTS_DIR / "Model_Comparison_Results.pptx"

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
COLOR_DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_SECTION_BG = RGBColor(0x16, 0x21, 0x3E)
COLOR_SUBHEADER_BG = RGBColor(0x0F, 0x3D, 0x3E)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)
COLOR_ACCENT = RGBColor(0x00, 0xD2, 0xFF)
COLOR_ACCENT2 = RGBColor(0x6C, 0x63, 0xFF)
COLOR_SUBTITLE = RGBColor(0xA0, 0xA0, 0xB0)
COLOR_GRAY = RGBColor(0x78, 0x78, 0x78)

# Comparison pairs: (section_title, model_a_dir, model_a_label, model_b_dir, model_b_label)
COMPARISON_PAIRS = [
    (
        "Fully Supervised Learning",
        "fsl_ST_UNet", "ST-UNet",
        "fsl_VSS_Net", "VSS-Net",
    ),
    (
        "Semi-Supervised Learning",
        "ssl_10_60_ite_3_student", "Student Iteration 3, L = 10, U = 60, no SDA",
        "ssl_10_60_SDA_ite_3_student", "Student Iteration 3, L = 10, U = 60, SDA",
    ),
    (
        "Weakly Supervised Learning",
        "wsl_RDFA_wsl_train_sscr_ablation", "RDFA SSCR",
        "wsl_SALE_wsl_train_sscr_ablation", "SALE SSCR",
    ),
]

# Selected cases: (category_dir, case_subpath, display_group)
SELECTED_CASES = [
    # Rat cases
    ("RatData_DIAS_Inference", "SRG_Angio/SRG_24_Cios/Aorta_1_dsa", "Rat Data"),
    ("RatData_DIAS_Inference", "SRG_Angio/SRG_29_Cios/Aorta_dsa", "Rat Data"),
    # Human cases
    ("HumanData_DIAS_Inference", "82A1/DSA/Segment_8_day_of_Tx_1", "Human Data"),
    ("HumanData_DIAS_Inference", "82A1/DSA/Left_Hepatic_Artery_Segment_4_Treatment_2/LHA_Treatment_2_Slightly_More_Proximal", "Human Data"),
]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def make_readable(name: str) -> str:
    """Convert folder names to readable titles."""
    return name.replace("_", " ").replace("&", " & ")


def get_summary_image_path(category_dir: str, model_dir: str, case_path: str) -> Path:
    """Get the path to a summary_comparison.png for a given model and case."""
    return RESULTS_DIR / category_dir / model_dir / case_path / "summary_comparison.png"


def split_summary_image(img_path: Path):
    """
    Split a summary_comparison.png into its left (input) and right (prediction) halves.

    The existing summary images are 2-panel composites: Input | Prediction.
    Returns (input_image, prediction_image) as PIL Images.
    """
    img = Image.open(img_path)
    w, h = img.size
    mid = w // 2
    input_img = img.crop((0, 0, mid, h))
    pred_img = img.crop((mid, 0, w, h))
    return input_img, pred_img


def build_comparison_image(
    input_img: Image.Image,
    pred_a_img: Image.Image,
    pred_b_img: Image.Image,
    label_input: str = "Input",
    label_a: str = "Model A",
    label_b: str = "Model B",
) -> Image.Image:
    """
    Build a 3-panel comparison image with labels above each panel.

    Layout: Input | Model A Prediction | Model B Prediction
    White background, black text labels above each panel.
    """
    # Resize all panels to the same height
    target_h = max(input_img.height, pred_a_img.height, pred_b_img.height)
    
    def resize_to_height(img, target_height):
        if img.height == target_height:
            return img
        ratio = target_height / img.height
        new_w = int(img.width * ratio)
        return img.resize((new_w, target_height), Image.LANCZOS)

    input_img = resize_to_height(input_img, target_h)
    pred_a_img = resize_to_height(pred_a_img, target_h)
    pred_b_img = resize_to_height(pred_b_img, target_h)

    # Label area dimensions
    label_height = 40
    gap = 10  # gap between panels
    total_w = input_img.width + pred_a_img.width + pred_b_img.width + 2 * gap
    total_h = target_h + label_height

    # Create canvas
    canvas = Image.new("RGB", (total_w, total_h), (255, 255, 255))
    draw = ImageDraw.Draw(canvas)

    # Try to load a reasonable font
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 24)
    except (IOError, OSError):
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf", 24)
        except (IOError, OSError):
            font = ImageFont.load_default()

    # Draw labels centered above each panel
    panels = [
        (0, input_img.width, label_input),
        (input_img.width + gap, pred_a_img.width, label_a),
        (input_img.width + gap + pred_a_img.width + gap, pred_b_img.width, label_b),
    ]

    for x_start, panel_w, label in panels:
        bbox = draw.textbbox((0, 0), label, font=font)
        text_w = bbox[2] - bbox[0]
        text_x = x_start + (panel_w - text_w) // 2
        text_y = (label_height - (bbox[3] - bbox[1])) // 2
        draw.text((text_x, text_y), label, fill=(0, 0, 0), font=font)

    # Paste panels
    canvas.paste(input_img, (0, label_height))
    canvas.paste(pred_a_img, (input_img.width + gap, label_height))
    canvas.paste(pred_b_img, (input_img.width + gap + pred_a_img.width + gap, label_height))

    return canvas


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_background(slide, color):
    """Set the slide background to a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs):
    """Add the opening title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide, COLOR_DARK_BG)

    # Title
    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(2.0), Inches(11.333), Inches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Model Comparison Results"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Accent line
    line = slide.shapes.add_shape(
        1, Inches(4.5), Inches(3.6), Inches(4.333), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    # Subtitle
    txBox2 = slide.shapes.add_textbox(
        Inches(1), Inches(4.0), Inches(11.333), Inches(1.0)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Side-by-Side Prediction Comparisons"
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLOR_SUBTITLE
    p2.alignment = PP_ALIGN.CENTER


def add_section_divider(prs, section_title: str, model_a_label: str, model_b_label: str):
    """Add a section divider slide for a paradigm (FSL/SSL/WSL)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLOR_SECTION_BG)

    # Accent bar on left
    bar = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.5), Inches(0.08), Inches(4.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()

    # Section title
    txBox = slide.shapes.add_textbox(
        Inches(1.2), Inches(2.0), Inches(10), Inches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.LEFT

    # Comparison info
    txBox2 = slide.shapes.add_textbox(
        Inches(1.2), Inches(3.8), Inches(10), Inches(1.0)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = f"{model_a_label}  vs  {model_b_label}"
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLOR_SUBTITLE
    p2.alignment = PP_ALIGN.LEFT


def add_subheader_slide(prs, section_title: str, group_name: str):
    """Add a sub-header slide (Rat Data / Human Data)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLOR_SUBHEADER_BG)

    # Section label (small, top)
    txBox0 = slide.shapes.add_textbox(
        Inches(1.2), Inches(1.8), Inches(10), Inches(0.6)
    )
    tf0 = txBox0.text_frame
    tf0.word_wrap = True
    p0 = tf0.paragraphs[0]
    p0.text = section_title
    p0.font.size = Pt(16)
    p0.font.color.rgb = COLOR_ACCENT
    p0.alignment = PP_ALIGN.LEFT

    # Group title
    txBox = slide.shapes.add_textbox(
        Inches(1.2), Inches(2.5), Inches(10), Inches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = group_name
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.LEFT

    # Accent line
    line = slide.shapes.add_shape(
        1, Inches(1.2), Inches(4.0), Inches(3.0), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT2
    line.line.fill.background()


def add_comparison_slide(
    prs,
    case_display_name: str,
    comparison_subtitle: str,
    comparison_image_path: str,
):
    """Add a content slide with the 3-panel comparison image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLOR_WHITE)

    # Title: case name
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(12), Inches(0.7)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = case_display_name
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLACK
    p.alignment = PP_ALIGN.LEFT

    # Subtitle: comparison pair
    txBox2 = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.85), Inches(12), Inches(0.5)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = comparison_subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_GRAY
    p2.alignment = PP_ALIGN.LEFT

    # Add the comparison image
    img_top = Inches(1.5)
    img_max_width = SLIDE_WIDTH - Inches(1.0)
    img_max_height = SLIDE_HEIGHT - img_top - Inches(0.3)

    try:
        with Image.open(comparison_image_path) as img:
            img_w, img_h = img.size

        # Calculate scale to fit within bounds
        scale_w = img_max_width / Emu(int(img_w * 914400 / 96))
        scale_h = img_max_height / Emu(int(img_h * 914400 / 96))
        scale = min(scale_w, scale_h, 1.0)

        final_w = Emu(int(img_w * 914400 / 96 * scale))
        final_h = Emu(int(img_h * 914400 / 96 * scale))

        # Center horizontally
        left = (SLIDE_WIDTH - final_w) // 2
        slide.shapes.add_picture(comparison_image_path, left, img_top, final_w, final_h)
    except Exception as e:
        print(f"  ERROR loading image {comparison_image_path}: {e}")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def case_display_name(case_path: str) -> str:
    """Convert a case subpath into a human-readable title."""
    return make_readable(case_path.replace("/", " / "))


def main():
    print(f"Results directory: {RESULTS_DIR}")
    print(f"Output file: {OUTPUT_FILE}\n")

    # Create a temp directory for composite images
    tmp_dir = RESULTS_DIR / "_comparison_tmp"
    tmp_dir.mkdir(exist_ok=True)

    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    add_title_slide(prs)
    print("Added title slide")

    for section_title, model_a_dir, model_a_label, model_b_dir, model_b_label in COMPARISON_PAIRS:
        # Section divider
        add_section_divider(prs, section_title, model_a_label, model_b_label)
        print(f"\nSection: {section_title} ({model_a_label} vs {model_b_label})")

        # Group cases by display group (Rat Data, Human Data)
        current_group = None

        for category_dir, case_path, display_group in SELECTED_CASES:
            # Add sub-header if group changed
            if display_group != current_group:
                add_subheader_slide(prs, section_title, display_group)
                current_group = display_group
                print(f"  Sub-header: {display_group}")

            # Find the summary images for both models
            img_a_path = get_summary_image_path(category_dir, model_a_dir, case_path)
            img_b_path = get_summary_image_path(category_dir, model_b_dir, case_path)

            if not img_a_path.exists():
                print(f"  WARNING: Missing {model_a_label} image: {img_a_path}")
                continue
            if not img_b_path.exists():
                print(f"  WARNING: Missing {model_b_label} image: {img_b_path}")
                continue

            # Extract panels
            input_img, pred_a_img = split_summary_image(img_a_path)
            _, pred_b_img = split_summary_image(img_b_path)

            # Build composite
            composite = build_comparison_image(
                input_img, pred_a_img, pred_b_img,
                label_input="Input",
                label_a=model_a_label,
                label_b=model_b_label,
            )

            # Save composite to temp
            safe_name = f"{model_a_dir}_vs_{model_b_dir}_{case_path.replace('/', '_')}.png"
            composite_path = tmp_dir / safe_name
            composite.save(str(composite_path))

            # Add slide
            display_name = case_display_name(case_path)
            subtitle = f"{model_a_label}  vs  {model_b_label}"
            add_comparison_slide(prs, display_name, subtitle, str(composite_path))
            print(f"  Slide: {display_name}")

    # Save
    prs.save(str(OUTPUT_FILE))
    total_slides = len(prs.slides)
    print(f"\nDone! Saved {total_slides} slides to:\n  {OUTPUT_FILE}")

    # Clean up temp images
    import shutil
    shutil.rmtree(tmp_dir, ignore_errors=True)
    print("Cleaned up temporary files.")


if __name__ == "__main__":
    main()
